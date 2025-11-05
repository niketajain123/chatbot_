import os
import json
from datetime import datetime
from PIL import Image
import PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation

class FileService:
    def __init__(self, upload_dir='uploads'):
        self.upload_dir = upload_dir
        os.makedirs(upload_dir, exist_ok=True)
    
    def save_file(self, file):
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{timestamp}_{file.filename}"
        filepath = os.path.join(self.upload_dir, filename)
        file.save(filepath)
        
        file_info = {
            'original_name': file.filename,
            'saved_name': filename,
            'filepath': filepath,
            'size': os.path.getsize(filepath),
            'upload_time': datetime.now().isoformat(),
            'type': self.get_file_type(file.filename)
        }
        return file_info
    
    def get_file_type(self, filename):
        ext = filename.lower().split('.')[-1]
        type_map = {
            'txt': 'text', 'md': 'text',
            'pdf': 'pdf',
            'docx': 'word', 'doc': 'word',
            'xlsx': 'excel', 'xls': 'excel',
            'pptx': 'powerpoint', 'ppt': 'powerpoint',
            'jpg': 'image', 'jpeg': 'image', 'png': 'image', 'gif': 'image'
        }
        return type_map.get(ext, 'unknown')
    
    def extract_text(self, filepath, file_type):
        try:
            if file_type == 'text':
                with open(filepath, 'r', encoding='utf-8') as f:
                    return f.read()
            
            elif file_type == 'pdf':
                with open(filepath, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    text = ""
                    for page in reader.pages:
                        text += page.extract_text() + "\n"
                    return text
            
            elif file_type == 'word':
                doc = Document(filepath)
                return "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            elif file_type == 'excel':
                wb = openpyxl.load_workbook(filepath)
                text = ""
                for sheet in wb.worksheets:
                    text += f"Sheet: {sheet.title}\n"
                    for row in sheet.iter_rows(values_only=True):
                        text += "\t".join([str(cell) if cell else "" for cell in row]) + "\n"
                return text
            
            elif file_type == 'powerpoint':
                prs = Presentation(filepath)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            
            elif file_type == 'image':
                return f"[Image: {os.path.basename(filepath)}]"
            
            return ""
        except Exception as e:
            return f"Error extracting text: {str(e)}"
    
    def get_file_info(self, filepath):
        if not os.path.exists(filepath):
            return None
        
        stat = os.stat(filepath)
        return {
            'name': os.path.basename(filepath),
            'size': stat.st_size,
            'modified': datetime.fromtimestamp(stat.st_mtime).isoformat(),
            'type': self.get_file_type(filepath)
        }
    
    def list_files(self):
        files = []
        for filename in os.listdir(self.upload_dir):
            filepath = os.path.join(self.upload_dir, filename)
            if os.path.isfile(filepath):
                files.append(self.get_file_info(filepath))
        return sorted(files, key=lambda x: x['modified'], reverse=True)
    
    def delete_file(self, filename):
        filepath = os.path.join(self.upload_dir, filename)
        if os.path.exists(filepath):
            os.remove(filepath)
            return True
        return False

file_service = FileService()